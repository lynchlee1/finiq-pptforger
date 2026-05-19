import json
import os
import sys
from pptx import Presentation
from jsonschema import validate, ValidationError

class PPTGenerator:
    def __init__(self, schema_path):
        with open(schema_path, 'r', encoding='utf-8') as f:
            self.schema = json.load(f)

    def validate_json(self, data):
        try:
            validate(instance=data, schema=self.schema)
            return True, ""
        except ValidationError as e:
            return False, str(e)

    def generate(self, input_json_path, output_pptx_path=None):
        try:
            with open(input_json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
        except Exception as e:
            print(f"Error loading JSON: {e}")
            return False

        is_valid, error_msg = self.validate_json(data)
        if not is_valid:
            print(f"JSON Validation Error: {error_msg}")
            return False

        template_path = data['template']
        if not os.path.exists(template_path):
            print(f"Template not found: {template_path}")
            return False

        if output_pptx_path is None:
            output_pptx_path = f"generated_{os.path.basename(template_path)}"

        try:
            prs = Presentation(template_path)
            
            for slide_data in data['slides']:
                slide_index = slide_data['slide_index']
                if slide_index < 0 or slide_index >= len(prs.slides):
                    print(f"Warning: slide_index {slide_index} out of range (0-{len(prs.slides)-1}). Skipping.")
                    continue
                
                slide = prs.slides[slide_index]
                self._process_slide(slide, slide_data)

            prs.save(output_pptx_path)
            print(f"Successfully saved PPTX to {output_pptx_path}")
            return True
        except Exception as e:
            print(f"Error during PPTX generation: {e}")
            return False

    def _process_slide(self, slide, slide_data):
        for key, value in slide_data.items():
            if key == 'slide_index':
                continue
            
            placeholder = f"{{{{{key}}}}}"
            self._replace_text_in_shapes(slide.shapes, placeholder, str(value))

    def _replace_text_in_shapes(self, shapes, placeholder, value):
        for shape in shapes:
            # 1. Handle shapes with text frames (Text boxes, AutoShapes)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)
            
            # 2. Handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, value)
            
            # 3. Handle group shapes (recursive)
            if shape.shape_type == 6: # Group shape
                self._replace_text_in_shapes(shape.shapes, placeholder, value)

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 src/ppt_generator.py <input_json> [output_pptx]")
        sys.exit(1)

    input_json = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else None
    
    # Path to schema.json relative to this script
    current_dir = os.path.dirname(os.path.abspath(__file__))
    schema_file = os.path.join(current_dir, "schema.json")
    
    generator = PPTGenerator(schema_file)
    generator.generate(input_json, output_pptx)

if __name__ == "__main__":
    main()
