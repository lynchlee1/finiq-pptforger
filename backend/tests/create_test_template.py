from pptx import Presentation
from pptx.util import Inches

def create_template(path):
    prs = Presentation()
    
    # Slide 0: Title and Subtitle
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Project: {{project_name}}"
    subtitle.text = "Author: {{author_name}}\nDate: {{date}}"
    
    # Slide 1: Table
    slide_layout = prs.slide_layouts[5] # Blank or Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Data Table"
    
    rows, cols = 2, 2
    left, top, width, height = Inches(2), Inches(2), Inches(4), Inches(1.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "{{cat1}}"
    table.cell(1, 1).text = "{{val1}}"
    
    prs.save(path)
    print(f"Test template created at {path}")

if __name__ == "__main__":
    create_template("template_ppt.pptx")
