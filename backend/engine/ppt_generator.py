import json
import os
import sys
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from jsonschema import validate, ValidationError

ALIGN_MAP = {
    PP_ALIGN.LEFT: "LEFT",
    PP_ALIGN.CENTER: "CENTER",
    PP_ALIGN.RIGHT: "RIGHT",
    PP_ALIGN.JUSTIFY: "JUSTIFY"
}

DEFAULT_THEME_COLORS = {
    MSO_THEME_COLOR.BACKGROUND_1: "#FFFFFF",
    MSO_THEME_COLOR.TEXT_1: "#000000",
    MSO_THEME_COLOR.BACKGROUND_2: "#E7E6E6",
    MSO_THEME_COLOR.TEXT_2: "#44546A",
    MSO_THEME_COLOR.ACCENT_1: "#4472C4",
    MSO_THEME_COLOR.ACCENT_2: "#ED7D31",
    MSO_THEME_COLOR.ACCENT_3: "#A5A5A5",
    MSO_THEME_COLOR.ACCENT_4: "#FFC000",
    MSO_THEME_COLOR.ACCENT_5: "#5B9BD5",
    MSO_THEME_COLOR.ACCENT_6: "#70AD47",
    MSO_THEME_COLOR.HYPERLINK: "#0563C1",
    MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "#954F72"
}

class PPTGenerator:
    def __init__(self, schema_path):
        with open(schema_path, 'r', encoding='utf-8') as f:
            self.schema = json.load(f)
        self.theme_colors = DEFAULT_THEME_COLORS.copy()

    def _extract_theme_colors(self, prs):
        """Extracts actual theme colors from the PPTX file."""
        self.theme_colors = DEFAULT_THEME_COLORS.copy()
        try:
            theme_part = None
            for rel in prs.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel.target_part
                    break
            if not theme_part and prs.slide_masters:
                for rel in prs.slide_masters[0].part.rels.values():
                    if "theme" in rel.reltype:
                        theme_part = rel.target_part
                        break
            if not theme_part: return

            theme_xml = theme_part.blob
            root = ET.fromstring(theme_xml)
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            clr_scheme = root.find('.//a:clrScheme', ns)
            if clr_scheme is not None:
                mapping = {
                    'dk1': MSO_THEME_COLOR.TEXT_1, 'lt1': MSO_THEME_COLOR.BACKGROUND_1,
                    'dk2': MSO_THEME_COLOR.TEXT_2, 'lt2': MSO_THEME_COLOR.BACKGROUND_2,
                    'accent1': MSO_THEME_COLOR.ACCENT_1, 'accent2': MSO_THEME_COLOR.ACCENT_2,
                    'accent3': MSO_THEME_COLOR.ACCENT_3, 'accent4': MSO_THEME_COLOR.ACCENT_4,
                    'accent5': MSO_THEME_COLOR.ACCENT_5, 'accent6': MSO_THEME_COLOR.ACCENT_6,
                    'hlink': MSO_THEME_COLOR.HYPERLINK, 'folHlink': MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
                }
                for color_node in clr_scheme:
                    name = color_node.tag.split('}')[-1]
                    if name in mapping:
                        srgb_clr = color_node.find('.//a:srgbClr', ns)
                        if srgb_clr is not None:
                            val = srgb_clr.get('val')
                            self.theme_colors[mapping[name]] = f"#{val}"
                        else:
                            sys_clr = color_node.find('.//a:sysClr', ns)
                            if sys_clr is not None and sys_clr.get('lastClr'):
                                self.theme_colors[mapping[name]] = f"#{sys_clr.get('lastClr')}"
        except Exception as e:
            sys.stderr.write(f"Warning: Failed to extract theme colors: {e}\n")

    def _apply_brightness(self, hex_color, brightness):
        try:
            if not hex_color or brightness == 0: return hex_color
            color = hex_color.lstrip('#')
            if len(color) != 6: return hex_color
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            if brightness > 0:
                r = int(r + (255 - r) * brightness)
                g = int(g + (255 - g) * brightness)
                b = int(b + (255 - b) * brightness)
            else:
                r = int(r * (1 + brightness))
                g = int(g * (1 + brightness))
                b = int(b * (1 + brightness))
            return f"#{max(0, min(255, r)):02X}{max(0, min(255, g)):02X}{max(0, min(255, b)):02X}"
        except: return hex_color

    def _get_color_hex(self, color_obj):
        try:
            hex_val = None
            if color_obj.type == MSO_COLOR_TYPE.RGB:
                hex_val = f"#{str(color_obj.rgb)}"
            elif color_obj.type == MSO_COLOR_TYPE.SCHEME:
                hex_val = self.theme_colors.get(color_obj.theme_color)
            if hex_val and hasattr(color_obj, 'brightness') and color_obj.brightness != 0:
                hex_val = self._apply_brightness(hex_val, color_obj.brightness)
            return hex_val
        except: return None

    def _flatten_color(self, color_obj):
        try:
            if color_obj.type == MSO_COLOR_TYPE.SCHEME or (
                color_obj.type == MSO_COLOR_TYPE.RGB and 
                hasattr(color_obj, 'brightness') and color_obj.brightness != 0
            ):
                hex_val = self._get_color_hex(color_obj)
                if hex_val:
                    r, g, b = int(hex_val[1:3], 16), int(hex_val[3:5], 16), int(hex_val[5:7], 16)
                    color_obj.rgb = RGBColor(r, g, b)
                    if hasattr(color_obj, 'brightness'): color_obj.brightness = 0
        except: pass

    def _flatten_all_colors(self, prs):
        for master in prs.slide_masters:
            self._flatten_shapes_recursive(master.shapes)
            for layout in master.slide_layouts: self._flatten_shapes_recursive(layout.shapes)
        for slide in prs.slides:
            self._flatten_shapes_recursive(slide.shapes)
            try:
                if slide.background.fill.type == MSO_FILL.SOLID:
                    self._flatten_color(slide.background.fill.fore_color)
            except: pass

    def _flatten_shapes_recursive(self, shapes):
        for shape in shapes:
            try:
                if hasattr(shape, 'fill') and shape.fill.type == MSO_FILL.SOLID:
                    self._flatten_color(shape.fill.fore_color)
            except: pass
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs: self._flatten_color(run.font.color)
            if shape.has_table:
                self._style_table(shape.table)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._flatten_shapes_recursive(shape.shapes)

    def _get_fill_color(self, fill):
        try:
            if fill.type == MSO_FILL.SOLID: return self._get_color_hex(fill.fore_color)
            return None
        except: return None

    def validate_json(self, data):
        try:
            validate(instance=data, schema=self.schema)
            return True, ""
        except ValidationError as e: return False, str(e)

    def _resolve_template_path(self, path):
        if os.path.isdir(path):
            pptx_files = [f for f in os.listdir(path) if f.endswith('.pptx') and not f.startswith('~$')]
            if not pptx_files: return path
            dir_name = os.path.basename(path.rstrip(os.sep))
            for f in pptx_files:
                if os.path.splitext(f)[0].lower() == dir_name.lower(): return os.path.join(path, f)
            for f in pptx_files:
                if f.lower() == 'deal-summary.pptx': return os.path.join(path, f)
            return os.path.join(path, pptx_files[0])
        return path

    def generate(self, input_json_path, output_pptx_path=None):
        try:
            with open(input_json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
        except Exception as e:
            sys.stderr.write(f"Error loading JSON: {e}\n")
            return False

        is_valid, error_msg = self.validate_json(data)
        if not is_valid:
            sys.stderr.write(f"JSON Validation Error: {error_msg}\n")
            return False

        template_path = self._resolve_template_path(data['template'])
        if not os.path.exists(template_path):
            sys.stderr.write(f"Template not found: {template_path}\n")
            return False

        if output_pptx_path is None:
            output_pptx_path = f"generated_{os.path.basename(template_path)}"

        try:
            prs = Presentation(template_path)
            self._extract_theme_colors(prs)
            
            for slide_data in data['slides']:
                slide_index = slide_data['slide_index']
                if slide_index < 0 or slide_index >= len(prs.slides): continue
                slide = prs.slides[slide_index]
                self._process_slide(slide, slide_data)

            self._flatten_all_colors(prs)
            prs.save(output_pptx_path)
            print(json.dumps({"success": True, "path": os.path.abspath(output_pptx_path)}))
            return True
        except Exception as e:
            sys.stderr.write(f"Error during PPTX generation: {e}\n")
            return False

    def _process_slide(self, slide, slide_data):
        for key, value in slide_data.items():
            if key == 'slide_index': continue
            placeholder = f"{{{{{key}}}}}"
            self._replace_text_in_shapes(slide.shapes, placeholder, str(value))
        self._apply_global_styling(slide.shapes)

    def _apply_global_styling(self, shapes):
        for shape in shapes:
            if shape.has_table: self._style_table(shape.table)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP: self._apply_global_styling(shape.shapes)

    def _set_font_all_languages(self, font, font_name):
        try:
            font.name = font_name
            rPr = font._element
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                from pptx.oxml.xmlchemy import OxmlElement
                ea = OxmlElement('a:ea')
                rPr.insert(0, ea)
            ea.set('typeface', font_name)
        except: pass

    def _style_table(self, table):
        try: self._fix_table_borders(table)
        except: pass

        for row in table.rows:
            for cell in row.cells:
                # Use original margins from template instead of forcing 0.05"
                try:
                    if cell.fill.type == MSO_FILL.SOLID: self._flatten_color(cell.fill.fore_color)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                except: pass
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs: self._flatten_color(run.font.color)

    def _fix_table_borders(self, table):
        for row in table.rows:
            for cell in row.cells:
                tcPr = cell._tc.get_or_add_tcPr()
                for side in ['lnT', 'lnB', 'lnL', 'lnR']:
                    ln = tcPr.find(qn(f'a:{side}'))
                    if ln is None:
                        from pptx.oxml.xmlchemy import OxmlElement
                        ln = OxmlElement(f'a:{side}')
                        ln.set('w', '12700')
                        tcPr.append(ln)
                    solidFill = ln.find(qn('a:solidFill'))
                    if solidFill is not None: ln.remove(solidFill)
                    from pptx.oxml.xmlchemy import OxmlElement
                    solidFill = OxmlElement('a:solidFill')
                    srgbClr = OxmlElement('a:srgbClr')
                    srgbClr.set('val', 'BFBFBF')
                    solidFill.append(srgbClr)
                    ln.append(solidFill)

    def _replace_text_in_shapes(self, shapes, placeholder, value):
        for shape in shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            orig_font_name = run.font.name
                            orig_font_size = run.font.size
                            orig_bold = run.font.bold
                            run.text = run.text.replace(placeholder, value)
                            if orig_font_name: self._set_font_all_languages(run.font, orig_font_name)
                            if orig_font_size: run.font.size = orig_font_size
                            if orig_bold is not None: run.font.bold = orig_bold
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if placeholder in run.text:
                                        orig_font_name = run.font.name
                                        orig_font_size = run.font.size
                                        orig_bold = run.font.bold
                                        run.text = run.text.replace(placeholder, value)
                                        if orig_font_name: self._set_font_all_languages(run.font, orig_font_name)
                                        if orig_font_size: run.font.size = orig_font_size
                                        if orig_bold is not None: run.font.bold = orig_bold
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._replace_text_in_shapes(shape.shapes, placeholder, value)

    def scan_template(self, template_path):
        template_path = self._resolve_template_path(template_path)
        if not os.path.exists(template_path): return False
        import re
        pattern = re.compile(r'\{\{([^}]+)\}\}')
        slides_info = []
        try:
            prs = Presentation(template_path)
            self._extract_theme_colors(prs)
            slide_width, slide_height = prs.slide_width, prs.slide_height
            for i, slide in enumerate(prs.slides):
                elements = []
                bg_color = self._get_fill_color(slide.background.fill) if hasattr(slide, 'background') else None
                try: self._extract_geometry_from_shapes(slide.slide_layout.slide_master.shapes, pattern, elements)
                except: pass
                try: self._extract_geometry_from_shapes(slide.slide_layout.shapes, pattern, elements)
                except: pass
                self._extract_geometry_from_shapes(slide.shapes, pattern, elements)
                slide_keys = {p for el in elements if el.get('placeholders') for p in el['placeholders']}
                slides_info.append({"slide_index": i, "background_color": bg_color, "keys": list(slide_keys), "elements": elements})
            print(json.dumps({"slide_count": len(prs.slides), "slide_width": slide_width, "slide_height": slide_height, "slides": slides_info}))
            return True
        except: return False

    def _extract_geometry_from_shapes(self, shapes, pattern, elements):
        for shape in shapes:
            try: left, top, width, height = shape.left, shape.top, shape.width, shape.height
            except: continue
            fill_color = self._get_fill_color(shape.fill) if hasattr(shape, 'fill') else None
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text: continue
                    font_size, text_color, is_bold = None, None, False
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        font_size, text_color, is_bold = run.font.size, self._get_color_hex(run.font.color), run.font.bold
                    matches = list(pattern.finditer(text))
                    el = {"left": left, "top": top, "width": width, "height": height, "font_size": font_size, "text_color": text_color, "fill_color": fill_color, "is_bold": is_bold, "alignment": ALIGN_MAP.get(paragraph.alignment, "LEFT") if paragraph.alignment else "LEFT"}
                    if matches:
                        el["type"], el["original_text"], el["placeholders"] = "placeholder_container", text, [m.group(1) for m in matches]
                    else:
                        el["type"], el["text"] = "static", text
                    elements.append(el)
            elif shape.has_table:
                table_el = {"type": "table", "left": left, "top": top, "width": width, "height": height, "rows": []}
                for r_idx, row in enumerate(shape.table.rows):
                    row_data = {"height": row.height, "cells": []}
                    for c_idx, cell in enumerate(row.cells):
                        cell_info = {"left": left, "top": top, "width": shape.table.columns[c_idx].width, "height": row.height, "text": cell.text.strip(), "fill_color": self._get_fill_color(cell.fill) if hasattr(cell, 'fill') else None}
                        matches = list(pattern.finditer(cell.text))
                        if matches:
                            cell_info["type"], cell_info["placeholders"] = "placeholder_container", [m.group(1) for m in matches]
                        else: cell_info["type"] = "static"
                        row_data["cells"].append(cell_info)
                    table_el["rows"].append(row_data)
                elements.append(table_el)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP: self._extract_geometry_from_shapes(shape.shapes, pattern, elements)
            elif shape.shape_type == 9: elements.append({"type": "line", "left": left, "top": top, "width": width, "height": height, "stroke_color": self._get_color_hex(shape.line.color) if hasattr(shape, 'line') else "#000000"})
            elif fill_color: elements.append({"type": "shape", "left": left, "top": top, "width": width, "height": height, "fill_color": fill_color})

    def move_slide_to_front(self, input_pptx_path, slide_index, output_pptx_path):
        try:
            prs = Presentation(input_pptx_path)
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise ValueError(f"Invalid slide index: {slide_index}")
            slide_ids = prs.slides._sldIdLst
            target = list(slide_ids)[slide_index]
            slide_ids.remove(target)
            slide_ids.insert(0, target)
            self._force_preview_slide_fonts(prs.slides[0])
            self._lock_preview_table_geometry(prs.slides[0])
            prs.save(output_pptx_path)
            print(json.dumps({"success": True, "path": os.path.abspath(output_pptx_path)}))
            return True
        except Exception as e:
            sys.stderr.write(f"Error during slide reorder: {e}\n")
            return False

    def postprocess_preview_image(self, pptx_path, image_path, output_path):
        try:
            from PIL import Image

            prs = Presentation(pptx_path)
            slide = prs.slides[0]
            image = Image.open(image_path).convert("RGB")
            image = self._replace_quicklook_background(image)

            table_shapes = [
                shape for shape in slide.shapes
                if getattr(shape, "has_table", False) and len(shape.table.rows) > 0
            ]
            for shape in sorted(table_shapes, key=lambda s: s.top):
                image = self._redraw_preview_table(image, shape, prs.slide_width, prs.slide_height)

            image.save(output_path)
            print(json.dumps({"success": True, "path": os.path.abspath(output_path)}))
            return True
        except Exception as e:
            sys.stderr.write(f"Error during preview image postprocess: {e}\n")
            return False

    def _replace_quicklook_background(self, image):
        pixels = image.load()
        width, height = image.size
        for y in range(height):
            for x in range(width):
                r, g, b = pixels[x, y]
                if abs(r - 172) <= 8 and abs(g - 178) <= 8 and abs(b - 187) <= 8:
                    pixels[x, y] = (255, 255, 255)
        return image

    def _shape_pixel_box(self, shape, slide_width, slide_height, image_size):
        image_width, image_height = image_size
        left = round(shape.left / slide_width * image_width)
        top = round(shape.top / slide_height * image_height)
        right = round((shape.left + shape.width) / slide_width * image_width)
        bottom = round((shape.top + shape.height) / slide_height * image_height)
        return left, top, right, bottom

    def _redraw_preview_table(self, image, shape, slide_width, slide_height):
        from PIL import Image, ImageDraw, ImageFont

        left, top, right, bottom = self._shape_pixel_box(shape, slide_width, slide_height, image.size)
        rendered_lines = self._find_table_horizontal_lines(image, left, right, top, bottom, 1000)
        rendered_bottom = rendered_lines[-1][1] + 3 if len(rendered_lines) >= len(shape.table.rows) + 1 else bottom
        if rendered_bottom > bottom + 3:
            fixed = Image.new("RGB", image.size, "white")
            fixed.paste(image.crop((0, 0, image.size[0], top)), (0, 0))
            fixed.paste(image.crop((0, rendered_bottom, image.size[0], image.size[1])), (0, bottom))
            image = fixed

        draw = ImageDraw.Draw(image)
        draw.rectangle((max(0, left - 4), max(0, top - 4), min(image.size[0], right + 32), bottom + 4), fill=(255, 255, 255))

        scale_x = image.size[0] / slide_width
        scale_y = image.size[1] / slide_height
        row_tops = [top]
        for row in shape.table.rows:
            row_tops.append(row_tops[-1] + round(row.height * scale_y))
        row_tops[-1] = bottom

        col_lefts = [left]
        for col in shape.table.columns:
            col_lefts.append(col_lefts[-1] + round(col.width * scale_x))
        col_lefts[-1] = right

        for r_idx, row in enumerate(shape.table.rows):
            for c_idx, cell in enumerate(row.cells):
                cell_box = (col_lefts[c_idx], row_tops[r_idx], col_lefts[c_idx + 1], row_tops[r_idx + 1])
                fill = self._preview_cell_fill(cell)
                draw.rectangle(cell_box, fill=fill)
                self._draw_preview_cell_text(draw, cell, cell_box, scale_x, scale_y)

        for r_idx, row in enumerate(shape.table.rows):
            for c_idx, cell in enumerate(row.cells):
                cell_box = (col_lefts[c_idx], row_tops[r_idx], col_lefts[c_idx + 1], row_tops[r_idx + 1])
                self._draw_preview_cell_borders(draw, cell, cell_box, scale_x)
        return image

    def _preview_cell_fill(self, cell):
        try:
            solid = cell._tc.get_or_add_tcPr().find(qn('a:solidFill'))
            if solid is not None:
                srgb = solid.find(qn('a:srgbClr'))
                if srgb is not None and srgb.get('val'):
                    return self._hex_to_rgb(f"#{srgb.get('val')}")
        except: pass
        return (255, 255, 255)

    def _preview_cell_margin(self, cell, name, default):
        try:
            value = cell._tc.get_or_add_tcPr().get(name)
            return int(value) if value is not None else default
        except: return default

    def _draw_preview_cell_text(self, draw, cell, cell_box, scale_x, scale_y):
        left, top, right, bottom = cell_box
        mar_l = round(self._preview_cell_margin(cell, 'marL', 91440) * scale_x)
        mar_r = round(self._preview_cell_margin(cell, 'marR', 91440) * scale_x)
        mar_t = round(self._preview_cell_margin(cell, 'marT', 18000) * scale_y)
        mar_b = round(self._preview_cell_margin(cell, 'marB', 18000) * scale_y)
        x = left + mar_l
        max_x = right - mar_r

        runs = []
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    runs.append(run)
        if not runs:
            return

        metrics = []
        total_width = 0
        max_height = 0
        for run in runs:
            font = self._preview_font(run)
            bbox = draw.textbbox((0, 0), run.text, font=font)
            width = bbox[2] - bbox[0]
            height = bbox[3] - bbox[1]
            metrics.append((run, font, width, height))
            total_width += width
            max_height = max(max_height, height)

        tcPr = cell._tc.get_or_add_tcPr()
        anchor = tcPr.get('anchor')
        if anchor == 'ctr':
            y = top + max(mar_t, ((bottom - top) - max_height) // 2 - 2)
        else:
            y = top + mar_t

        for run, font, width, _ in metrics:
            if x >= max_x:
                break
            color = self._preview_run_color(run)
            draw.text((x, y), run.text, font=font, fill=color)
            if self._preview_run_underline(run):
                underline_y = y + max_height + 2
                draw.line((x, underline_y, min(x + width, max_x), underline_y), fill=color, width=1)
            x += width

    def _preview_font(self, run):
        from PIL import ImageFont

        text = run.text or ""
        bold = bool(run.font.bold)
        if any('\uac00' <= ch <= '\ud7a3' for ch in text):
            path = "/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/malgunbd.ttf" if bold else "/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/malgun.ttf"
        else:
            path = "/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/Calibrib.ttf" if bold else "/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/Calibri.ttf"
        size_pt = run.font.size.pt if run.font.size else 11
        size_px = max(1, round(size_pt * 2400 / 9906000 * 914400 / 72))
        return ImageFont.truetype(path, size_px)

    def _preview_run_color(self, run):
        hex_color = self._get_color_hex(run.font.color)
        return self._hex_to_rgb(hex_color or "#000000")

    def _preview_run_underline(self, run):
        try: return bool(run.font.underline)
        except: return False

    def _draw_preview_cell_borders(self, draw, cell, cell_box, scale_x):
        left, top, right, bottom = cell_box
        for tag, coords in (
            ('a:lnL', (left, top, left, bottom)),
            ('a:lnR', (right, top, right, bottom)),
            ('a:lnT', (left, top, right, top)),
            ('a:lnB', (left, bottom, right, bottom)),
        ):
            line = cell._tc.get_or_add_tcPr().find(qn(tag))
            if line is None:
                color = (191, 191, 191)
                width = 1
            else:
                color = self._preview_line_color(line)
                width = max(1, round(int(line.get('w') or 12700) * scale_x))
            draw.line(coords, fill=color, width=width)

    def _preview_line_color(self, line):
        try:
            solid = line.find(qn('a:solidFill'))
            if solid is not None:
                srgb = solid.find(qn('a:srgbClr'))
                if srgb is not None and srgb.get('val'):
                    return self._hex_to_rgb(f"#{srgb.get('val')}")
        except: pass
        return (191, 191, 191)

    def _hex_to_rgb(self, hex_color):
        value = (hex_color or "#000000").lstrip("#")
        return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))

    def _find_table_horizontal_lines(self, image, left, right, top, bottom, row_count):
        width, height = image.size
        x0 = max(0, left)
        x1 = min(width, right)
        y0 = max(0, top - 12)
        y1 = min(height, bottom + max(bottom - top, 120))
        if x1 <= x0 or y1 <= y0:
            return []

        threshold = max(40, int((x1 - x0) * 0.35))
        rows = []
        pixels = image.load()
        for y in range(y0, y1):
            count = 0
            for x in range(x0, x1):
                r, g, b = pixels[x, y]
                if abs(r - g) < 10 and abs(g - b) < 10 and 110 <= r <= 215:
                    count += 1
            if count >= threshold:
                rows.append(y)

        groups = []
        for y in rows:
            if not groups or y > groups[-1][1] + 1:
                groups.append([y, y])
            else:
                groups[-1][1] = y
        return groups[:row_count + 1]

    def _force_preview_slide_fonts(self, slide):
        for shape in slide.shapes:
            self._force_preview_font_in_shape(shape)

    def _force_preview_font_in_shape(self, shape):
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    self._set_preview_font(run)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                self._set_preview_font(run)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._force_preview_font_in_shape(child)

    def _lock_preview_table_geometry(self, slide):
        for shape in slide.shapes:
            self._lock_table_geometry_in_shape(shape)

    def _lock_table_geometry_in_shape(self, shape):
        if shape.has_table and len(shape.table.rows) > 0:
            total_height = shape.height
            base_height = total_height // len(shape.table.rows)
            used_height = 0
            for idx, row in enumerate(shape.table.rows):
                if idx == len(shape.table.rows) - 1:
                    row.height = total_height - used_height
                else:
                    row.height = base_height
                    used_height += base_height
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._lock_table_geometry_in_shape(child)

    def _set_preview_font(self, run):
        try:
            rPr = run.font._element
            from pptx.oxml.xmlchemy import OxmlElement
            for tag, typeface in (
                ('a:latin', 'Calibri'),
                ('a:ea', 'Malgun Gothic'),
                ('a:cs', 'Malgun Gothic'),
            ):
                font_node = rPr.find(qn(tag))
                if font_node is None:
                    font_node = OxmlElement(tag)
                    rPr.insert(0, font_node)
                font_node.set('typeface', typeface)
        except: pass

def main():
    if len(sys.argv) < 2: sys.exit(1)
    command = sys.argv[1]
    current_dir = os.path.dirname(os.path.abspath(__file__))
    schema_file = os.path.join(current_dir, "../../shared/schema.json")
    generator = PPTGenerator(schema_file)
    if command == "generate": generator.generate(sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else None)
    elif command == "scan": generator.scan_template(sys.argv[2])
    elif command == "move-slide-to-front": generator.move_slide_to_front(sys.argv[2], int(sys.argv[3]), sys.argv[4])
    elif command == "postprocess-preview-image": generator.postprocess_preview_image(sys.argv[2], sys.argv[3], sys.argv[4])
    else: generator.generate(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)

if __name__ == "__main__": main()
