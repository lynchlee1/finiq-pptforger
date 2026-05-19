from pptx import Presentation

def inspect_pptx(path):
    prs = Presentation(path)
    print(f"Inspecting {path}:")
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    print(f"  Text: {paragraph.text}")
            if shape.has_table:
                for row in shape.table.rows:
                    cells_text = [cell.text for cell in row.cells]
                    print(f"  Table Row: {cells_text}")

if __name__ == "__main__":
    inspect_pptx("output.pptx")
